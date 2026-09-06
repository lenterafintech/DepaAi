"""Ekspor laporan hasil analisis ke berbagai format berkas.

Laporan disusun sekali menjadi rangkaian blok (judul, paragraf, poin, tabel), lalu
tiap format hanya menerjemahkan blok yang sama. Dengan begitu isi Word, PDF,
PowerPoint, Excel, HTML, dan Markdown tidak pernah berbeda satu sama lain.

Dua ragam laporan tersedia:

- **Ringkasan** — satu register pembaca saja, ringkas untuk dibagikan.
- **Laporan lengkap** — ketiga register, seluruh tabel, paragraf siap salin,
  rujukan, dan catatan metode.
"""

from __future__ import annotations

import io
import json
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd

from lentera_mva.narrative import AUDIENCE_LABELS, AUDIENCES, Laporan
from lentera_mva.report_html import laporan_html, laporan_html_semua
from lentera_mva.sintaks import bangkitkan

# Padanan aman untuk lambang yang tidak tersedia pada huruf bawaan PDF.
GANTI_LAMBANG = {
    "−": "-",
    "–": "-",
    "—": "-",
    "×": "x",
    "√": "akar ",
    "≥": ">=",
    "≤": "<=",
    "η": "eta",
    "χ": "chi",
    "β": "beta",
    "λ": "lambda",
    "α": "alpha",
    "²": "^2",
    "·": "-",
    "“": '"',
    "”": '"',
    "’": "'",
}
JALUR_FONT = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/DejaVuSans.ttf",
    "C:/Windows/Fonts/DejaVuSans.ttf",
]


@dataclass
class Blok:
    """Satu potongan isi laporan, bebas dari format berkasnya."""

    jenis: str  # judul | subjudul | paragraf | poin | tabel | meta | catatan
    teks: str = ""
    poin: list[str] = field(default_factory=list)
    tabel: pd.DataFrame | None = None
    catatan: str = ""


# --------------------------------------------------------------------------- #
# Penyusun isi
# --------------------------------------------------------------------------- #


def susun_blok(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> list[Blok]:
    """Ubah laporan menjadi rangkaian blok yang siap ditulis ke format apa pun."""
    pembaca_dipakai = list(AUDIENCES) if lengkap else [pembaca]
    judul = (
        "Laporan Lengkap Analisis Multivariat"
        if lengkap
        else f"Ringkasan {AUDIENCE_LABELS[pembaca]}"
    )
    blok: list[Blok] = [
        Blok("judul", judul),
        Blok(
            "meta",
            f"Sumber data: {laporan.dataset} · {laporan.n_baris} baris × "
            f"{laporan.n_kolom} kolom · {len(laporan.metode_terpakai)} metode · "
            f"{laporan.tanggal}",
        ),
        Blok("paragraf", laporan.headline),
        Blok("paragraf", laporan.subheadline),
    ]

    if laporan.lampu:
        blok.append(Blok("subjudul", "Status pemeriksaan"))
        blok.append(
            Blok(
                "tabel",
                tabel=pd.DataFrame(
                    {
                        "Pemeriksaan": [l.label for l in laporan.lampu],
                        "Nilai": [l.nilai for l in laporan.lampu],
                        "Status": [l.status_label for l in laporan.lampu],
                        "Catatan": [l.catatan for l in laporan.lampu],
                    }
                ),
            )
        )

    if laporan.pendorong:
        blok.append(Blok("subjudul", "Peringkat pendorong"))
        blok.append(Blok("catatan", catatan=f"Sumber: {laporan.pendorong_sumber}"))
        blok.append(
            Blok(
                "tabel",
                tabel=pd.DataFrame(
                    {
                        "Faktor": [p.nama for p in laporan.pendorong],
                        "Satuan": [p.satuan for p in laporan.pendorong],
                        "Nilai": [round(p.nilai, 4) for p in laporan.pendorong],
                        "Arah": [p.arah for p in laporan.pendorong],
                        "Signifikan": ["Ya" if p.signifikan else "Tidak" for p in laporan.pendorong],
                        "Catatan": [p.catatan for p in laporan.pendorong],
                    }
                ),
            )
        )

    for reg in pembaca_dipakai:
        tajuk = (
            f"Temuan — {AUDIENCE_LABELS[reg]}" if lengkap else "Temuan"
        )
        blok.append(Blok("subjudul", tajuk))
        for temuan in laporan.temuan:
            blok.append(Blok("paragraf", f"{temuan.judul} ({temuan.metode})"))
            blok.append(Blok("paragraf", temuan.teks(reg)))

    if lengkap or pembaca == "akademik":
        for nomor, (judul_tabel, tabel, catatan) in laporan.tabel.items():
            blok.append(Blok("subjudul", f"{nomor}. {judul_tabel}"))
            blok.append(Blok("tabel", tabel=tabel, catatan=catatan))

    if lengkap or pembaca == "akademik":
        if laporan.paragraf:
            blok.append(Blok("subjudul", "Kalimat siap salin"))
            for paragraf in laporan.paragraf:
                blok.append(Blok("paragraf", paragraf.bagian))
                blok.append(Blok("paragraf", paragraf.teks))

    if laporan.rekomendasi:
        blok.append(Blok("subjudul", "Rekomendasi tindakan"))
        blok.append(
            Blok(
                "poin",
                poin=[
                    f"{r.judul} (prioritas {r.prioritas}) — {r.alasan}"
                    for r in laporan.rekomendasi
                ],
            )
        )

    if laporan.keterbatasan:
        blok.append(Blok("subjudul", "Batas kesimpulan"))
        blok.append(Blok("poin", poin=list(laporan.keterbatasan)))

    if lengkap or pembaca == "akademik":
        if laporan.rujukan:
            blok.append(Blok("subjudul", "Rujukan ambang yang dipakai"))
            blok.append(Blok("poin", poin=list(laporan.rujukan)))

    if laporan.dilewati:
        blok.append(Blok("subjudul", "Analisis yang tidak dapat dijalankan"))
        blok.append(Blok("poin", poin=list(laporan.dilewati)))

    blok.append(
        Blok(
            "catatan",
            catatan=(
                "Disusun otomatis oleh Lentera MVA. Kesimpulan statistik menunjukkan pola "
                "dalam data, bukan bukti sebab-akibat; keputusan akhir tetap memerlukan "
                "pertimbangan konteks dan keahlian bidang."
            ),
        )
    )
    return blok


def _aman_pdf(teks: str) -> str:
    """Ganti lambang yang tidak ada pada huruf bawaan PDF dengan padanan ASCII."""
    hasil = str(teks)
    for lambang, ganti in GANTI_LAMBANG.items():
        hasil = hasil.replace(lambang, ganti)
    return hasil


# --------------------------------------------------------------------------- #
# Word (.docx)
# --------------------------------------------------------------------------- #


def ke_docx(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    dok = Document()
    gaya = dok.styles["Normal"]
    gaya.font.name = "Calibri"
    gaya.font.size = Pt(11)

    for blok in susun_blok(laporan, pembaca, lengkap):
        if blok.jenis == "judul":
            dok.add_heading(blok.teks, level=0)
        elif blok.jenis == "subjudul":
            dok.add_heading(blok.teks, level=1)
        elif blok.jenis == "meta":
            paragraf = dok.add_paragraph(blok.teks)
            paragraf.runs[0].font.size = Pt(9)
            paragraf.runs[0].font.color.rgb = RGBColor(0x6F, 0x7A, 0x91)
        elif blok.jenis == "paragraf":
            dok.add_paragraph(blok.teks)
        elif blok.jenis == "poin":
            for butir in blok.poin:
                dok.add_paragraph(butir, style="List Bullet")
        elif blok.jenis == "catatan" and blok.catatan:
            paragraf = dok.add_paragraph(blok.catatan)
            paragraf.runs[0].italic = True
            paragraf.runs[0].font.size = Pt(9)
        elif blok.jenis == "tabel" and blok.tabel is not None:
            tabel = blok.tabel.fillna("")
            objek = dok.add_table(rows=1, cols=len(tabel.columns))
            objek.style = "Light Grid Accent 1"
            for i, kolom in enumerate(tabel.columns):
                sel = objek.rows[0].cells[i]
                sel.text = str(kolom)
                for p in sel.paragraphs:
                    for r in p.runs:
                        r.bold = True
            for _, baris in tabel.iterrows():
                sel_baris = objek.add_row().cells
                for i, nilai in enumerate(baris):
                    sel_baris[i].text = str(nilai)
            if blok.catatan:
                catatan = dok.add_paragraph(f"Catatan. {blok.catatan}")
                catatan.runs[0].italic = True
                catatan.runs[0].font.size = Pt(9)

    penampung = io.BytesIO()
    dok.save(penampung)
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# Excel (.xlsx)
# --------------------------------------------------------------------------- #


def _nama_lembar(teks: str, dipakai: set[str]) -> str:
    """Nama lembar Excel: maksimal 31 karakter, tanpa lambang terlarang, dan unik."""
    bersih = "".join(c for c in str(teks) if c not in r"[]:*?/\\")[:31].strip() or "Lembar"
    calon = bersih
    urut = 2
    while calon.lower() in dipakai:
        akhiran = f" ({urut})"
        calon = bersih[: 31 - len(akhiran)] + akhiran
        urut += 1
    dipakai.add(calon.lower())
    return calon


def ke_xlsx(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    penampung = io.BytesIO()
    dipakai: set[str] = set()
    with pd.ExcelWriter(penampung, engine="openpyxl") as penulis:
        ringkas = pd.DataFrame(
            {
                "Keterangan": [
                    "Judul",
                    "Sumber data",
                    "Ukuran data",
                    "Metode dijalankan",
                    "Tanggal analisis",
                    "Kesimpulan utama",
                    "Penjelasan",
                ],
                "Isi": [
                    "Laporan lengkap" if lengkap else f"Ringkasan {AUDIENCE_LABELS[pembaca]}",
                    laporan.dataset,
                    f"{laporan.n_baris} baris × {laporan.n_kolom} kolom",
                    ", ".join(laporan.metode_terpakai),
                    laporan.tanggal,
                    laporan.headline,
                    laporan.subheadline,
                ],
            }
        )
        ringkas.to_excel(penulis, sheet_name=_nama_lembar("Ringkasan", dipakai), index=False)

        if laporan.lampu:
            pd.DataFrame(
                {
                    "Pemeriksaan": [l.label for l in laporan.lampu],
                    "Nilai": [l.nilai for l in laporan.lampu],
                    "Status": [l.status_label for l in laporan.lampu],
                    "Catatan": [l.catatan for l in laporan.lampu],
                }
            ).to_excel(penulis, sheet_name=_nama_lembar("Status", dipakai), index=False)

        if laporan.pendorong:
            pd.DataFrame(
                {
                    "Faktor": [p.nama for p in laporan.pendorong],
                    "Satuan": [p.satuan for p in laporan.pendorong],
                    "Nilai": [p.nilai for p in laporan.pendorong],
                    "Kekuatan relatif": [p.kekuatan for p in laporan.pendorong],
                    "Arah": [p.arah for p in laporan.pendorong],
                    "p-value": [p.p_value for p in laporan.pendorong],
                    "Catatan": [p.catatan for p in laporan.pendorong],
                }
            ).to_excel(penulis, sheet_name=_nama_lembar("Pendorong", dipakai), index=False)

        register = list(AUDIENCES) if lengkap else [pembaca]
        temuan = pd.DataFrame(
            {
                "Temuan": [t.judul for t in laporan.temuan],
                "Metode": [t.metode for t in laporan.temuan],
                "Ringkas": [t.ringkas for t in laporan.temuan],
                **{
                    AUDIENCE_LABELS[reg]: [t.teks(reg) for t in laporan.temuan]
                    for reg in register
                },
            }
        )
        temuan.to_excel(penulis, sheet_name=_nama_lembar("Temuan", dipakai), index=False)

        if laporan.rekomendasi:
            pd.DataFrame(
                {
                    "Rekomendasi": [r.judul for r in laporan.rekomendasi],
                    "Prioritas": [r.prioritas for r in laporan.rekomendasi],
                    "Alasan": [r.alasan for r in laporan.rekomendasi],
                }
            ).to_excel(penulis, sheet_name=_nama_lembar("Rekomendasi", dipakai), index=False)

        if laporan.keterbatasan:
            pd.DataFrame({"Batas kesimpulan": laporan.keterbatasan}).to_excel(
                penulis, sheet_name=_nama_lembar("Keterbatasan", dipakai), index=False
            )

        if lengkap or pembaca == "akademik":
            for nomor, (judul_tabel, tabel, _) in laporan.tabel.items():
                tabel.to_excel(
                    penulis, sheet_name=_nama_lembar(nomor, dipakai), index=False
                )
            if laporan.paragraf:
                pd.DataFrame(
                    {
                        "Bagian": [p.bagian for p in laporan.paragraf],
                        "Paragraf": [p.teks for p in laporan.paragraf],
                    }
                ).to_excel(penulis, sheet_name=_nama_lembar("Siap salin", dipakai), index=False)

    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #


def _daftarkan_font() -> str:
    """Pakai DejaVu bila tersedia agar lambang statistik tampil utuh."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for jalur in JALUR_FONT:
        berkas = Path(jalur)
        if berkas.exists():
            try:
                pdfmetrics.registerFont(TTFont("DejaVu", str(berkas)))
                tebal = berkas.with_name(berkas.stem + "-Bold" + berkas.suffix)
                if tebal.exists():
                    pdfmetrics.registerFont(TTFont("DejaVu-Bold", str(tebal)))
                    pdfmetrics.registerFontFamily("DejaVu", normal="DejaVu", bold="DejaVu-Bold")
                return "DejaVu"
            except Exception:  # noqa: BLE001 - huruf gagal dimuat, pakai bawaan
                break
    return "Helvetica"


def ke_pdf(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    font = _daftarkan_font()
    font_tebal = "DejaVu-Bold" if font == "DejaVu" else "Helvetica-Bold"
    perlu_sanitasi = font == "Helvetica"

    def teks(nilai: object) -> str:
        mentah = str(nilai)
        return _aman_pdf(mentah) if perlu_sanitasi else mentah

    dasar = getSampleStyleSheet()
    gaya_judul = ParagraphStyle(
        "JudulLentera", parent=dasar["Title"], fontName=font_tebal, fontSize=18,
        textColor=colors.HexColor("#131a2b"), alignment=TA_LEFT, spaceAfter=6,
    )
    gaya_sub = ParagraphStyle(
        "SubLentera", parent=dasar["Heading2"], fontName=font_tebal, fontSize=13,
        textColor=colors.HexColor("#26356b"), spaceBefore=14, spaceAfter=6,
    )
    gaya_isi = ParagraphStyle(
        "IsiLentera", parent=dasar["BodyText"], fontName=font, fontSize=10, leading=15,
        textColor=colors.HexColor("#1f2937"), spaceAfter=6,
    )
    gaya_meta = ParagraphStyle(
        "MetaLentera", parent=gaya_isi, fontSize=8.5,
        textColor=colors.HexColor("#6f7a91"),
    )

    penampung = io.BytesIO()
    dokumen = SimpleDocTemplate(
        penampung, pagesize=A4,
        leftMargin=20 * mm, rightMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
        title=("Laporan Lengkap" if lengkap else f"Ringkasan {AUDIENCE_LABELS[pembaca]}"),
        author="Lentera MVA",
    )
    isi: list = []
    for blok in susun_blok(laporan, pembaca, lengkap):
        if blok.jenis == "judul":
            isi.append(Paragraph(teks(blok.teks), gaya_judul))
        elif blok.jenis == "subjudul":
            isi.append(Paragraph(teks(blok.teks), gaya_sub))
        elif blok.jenis == "meta":
            isi.append(Paragraph(teks(blok.teks), gaya_meta))
        elif blok.jenis == "paragraf":
            isi.append(Paragraph(teks(blok.teks), gaya_isi))
        elif blok.jenis == "poin":
            for butir in blok.poin:
                isi.append(Paragraph("• " + teks(butir), gaya_isi))
        elif blok.jenis == "catatan" and blok.catatan:
            isi.append(Paragraph(teks(blok.catatan), gaya_meta))
        elif blok.jenis == "tabel" and blok.tabel is not None:
            tabel = blok.tabel.fillna("")
            data = [[Paragraph(f"<b>{teks(k)}</b>", gaya_meta) for k in tabel.columns]]
            for _, baris in tabel.iterrows():
                data.append([Paragraph(teks(v), gaya_meta) for v in baris])
            lebar = (dokumen.width) / max(len(tabel.columns), 1)
            objek = Table(data, colWidths=[lebar] * len(tabel.columns), repeatRows=1)
            objek.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eef1f8")),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#dde3ee")),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 4),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            isi.append(objek)
            if blok.catatan:
                isi.append(Spacer(1, 3))
                isi.append(Paragraph("Catatan. " + teks(blok.catatan), gaya_meta))
            isi.append(Spacer(1, 8))

    dokumen.build(isi)
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# PowerPoint (.pptx)
# --------------------------------------------------------------------------- #


def ke_pptx(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentasi = Presentation()
    presentasi.slide_width = Inches(13.333)
    presentasi.slide_height = Inches(7.5)
    tata_judul = presentasi.slide_layouts[0]
    tata_isi = presentasi.slide_layouts[1]
    tata_kosong = presentasi.slide_layouts[6]

    navy = RGBColor(0x26, 0x35, 0x6B)

    slide = presentasi.slides.add_slide(tata_judul)
    slide.shapes.title.text = laporan.headline
    slide.placeholders[1].text = (
        f"{laporan.subheadline}\n\n{laporan.dataset} · {laporan.n_baris} baris · "
        f"{laporan.tanggal}"
    )

    def slide_poin(judul: str, butir: list[str], maks: int = 6) -> None:
        """Satu slide berisi daftar poin; daftar panjang dipecah agar tetap terbaca."""
        for mulai in range(0, len(butir), maks):
            potongan = butir[mulai : mulai + maks]
            s = presentasi.slides.add_slide(tata_isi)
            s.shapes.title.text = judul if mulai == 0 else f"{judul} (lanjutan)"
            s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
            kerangka = s.placeholders[1].text_frame
            kerangka.clear()
            kerangka.word_wrap = True
            for i, isi in enumerate(potongan):
                paragraf = kerangka.paragraphs[0] if i == 0 else kerangka.add_paragraph()
                paragraf.text = isi
                paragraf.font.size = Pt(16)

    if laporan.lampu:
        slide_poin(
            "Status pemeriksaan",
            [f"{l.label}: {l.nilai} — {l.status_label}" for l in laporan.lampu],
        )
    if laporan.pendorong:
        slide_poin(
            "Peringkat pendorong",
            [
                f"{i}. {p.nama} ({p.satuan} {p.nilai:.3f}) — "
                f"{'signifikan' if p.signifikan else 'belum terbukti'}"
                for i, p in enumerate(laporan.pendorong[:8], start=1)
            ],
        )

    slide_poin("Poin kunci", [t.ringkas for t in laporan.temuan])

    register = pembaca if not lengkap else "eksekutif"
    for temuan in laporan.temuan:
        s = presentasi.slides.add_slide(tata_kosong)
        kotak_judul = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
        p = kotak_judul.text_frame.paragraphs[0]
        p.text = temuan.judul
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = navy

        kotak_isi = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5))
        kerangka = kotak_isi.text_frame
        kerangka.word_wrap = True
        kerangka.paragraphs[0].text = temuan.teks(register)
        kerangka.paragraphs[0].font.size = Pt(15)
        metode = kerangka.add_paragraph()
        metode.text = f"Metode: {temuan.metode}"
        metode.font.size = Pt(11)
        metode.font.color.rgb = RGBColor(0x6F, 0x7A, 0x91)

    if laporan.rekomendasi:
        slide_poin(
            "Rekomendasi tindakan",
            [f"{r.judul} — {r.alasan}" for r in laporan.rekomendasi],
            maks=4,
        )
    if laporan.keterbatasan:
        slide_poin("Batas kesimpulan", list(laporan.keterbatasan), maks=4)

    penampung = io.BytesIO()
    presentasi.save(penampung)
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# Format berbasis teks dan paket data
# --------------------------------------------------------------------------- #


def ke_html(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    berkas = laporan_html_semua(laporan) if lengkap else laporan_html(laporan, pembaca)
    return berkas.encode("utf-8")


def ke_markdown(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    if not lengkap:
        return laporan.markdown(pembaca).encode("utf-8")
    bagian = [laporan.markdown(reg) for reg in AUDIENCES]
    return ("\n\n---\n\n".join(bagian)).encode("utf-8")


def ke_json(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Isi laporan dalam bentuk terstruktur, untuk diolah kembali oleh sistem lain."""
    register = list(AUDIENCES) if lengkap else [pembaca]
    isi = {
        "dataset": laporan.dataset,
        "n_baris": laporan.n_baris,
        "n_kolom": laporan.n_kolom,
        "tanggal": laporan.tanggal,
        "metode": laporan.metode_terpakai,
        "headline": laporan.headline,
        "subheadline": laporan.subheadline,
        "lampu": [
            {"label": l.label, "nilai": l.nilai, "status": l.status, "catatan": l.catatan}
            for l in laporan.lampu
        ],
        "pendorong": [
            {
                "nama": p.nama,
                "satuan": p.satuan,
                "nilai": p.nilai,
                "kekuatan": p.kekuatan,
                "arah": p.arah,
                "p_value": None if pd.isna(p.p_value) else p.p_value,
                "signifikan": p.signifikan,
            }
            for p in laporan.pendorong
        ],
        "temuan": [
            {
                "judul": t.judul,
                "metode": t.metode,
                "ringkas": t.ringkas,
                **{reg: t.teks(reg) for reg in register},
            }
            for t in laporan.temuan
        ],
        "rekomendasi": [
            {"judul": r.judul, "prioritas": r.prioritas, "alasan": r.alasan}
            for r in laporan.rekomendasi
        ],
        "keterbatasan": laporan.keterbatasan,
        "dilewati": laporan.dilewati,
    }
    return json.dumps(isi, ensure_ascii=False, indent=2).encode("utf-8")


def ke_python(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Skrip Python yang menjalankan ulang analisis dengan pustaka yang sama."""
    return bangkitkan(laporan.konfig, "py").encode("utf-8")


def ke_r(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Skrip R sebagai pemeriksaan silang di luar Python."""
    return bangkitkan(laporan.konfig, "r").encode("utf-8")


def ke_zip(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Paket lengkap: laporan dalam beberapa format sekaligus, ditambah tabel CSV."""
    penampung = io.BytesIO()
    dasar = "laporan_lengkap" if lengkap else f"ringkasan_{pembaca}"
    with zipfile.ZipFile(penampung, "w", zipfile.ZIP_DEFLATED) as arsip:
        arsip.writestr(f"{dasar}.html", ke_html(laporan, pembaca, lengkap))
        arsip.writestr(f"{dasar}.md", ke_markdown(laporan, pembaca, lengkap))
        arsip.writestr(f"{dasar}.docx", ke_docx(laporan, pembaca, lengkap))
        arsip.writestr(f"{dasar}.xlsx", ke_xlsx(laporan, pembaca, lengkap))
        arsip.writestr(f"{dasar}.pdf", ke_pdf(laporan, pembaca, lengkap))
        arsip.writestr(f"{dasar}.json", ke_json(laporan, pembaca, lengkap))
        arsip.writestr("sintaks/analisis.py", ke_python(laporan, pembaca, lengkap))
        arsip.writestr("sintaks/analisis.R", ke_r(laporan, pembaca, lengkap))
        for nomor, (judul_tabel, tabel, _) in laporan.tabel.items():
            nama = nomor.lower().replace(" ", "_")
            arsip.writestr(f"tabel/{nama}.csv", tabel.to_csv(index=False))
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# Daftar format
# --------------------------------------------------------------------------- #


@dataclass(frozen=True)
class Format:
    kode: str
    nama: str
    ekstensi: str
    mime: str
    keterangan: str


FORMAT: dict[str, Format] = {
    "docx": Format("docx", "Word", "docx",
                   "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                   "Dokumen siap disunting dan dicetak."),
    "pdf": Format("pdf", "PDF", "pdf", "application/pdf",
                  "Tata letak tetap, cocok untuk lampiran resmi."),
    "xlsx": Format("xlsx", "Excel", "xlsx",
                   "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                   "Tiap bagian menjadi lembar tersendiri; angka siap diolah ulang."),
    "pptx": Format("pptx", "PowerPoint", "pptx",
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                   "Slide siap dipresentasikan."),
    "html": Format("html", "HTML", "html", "text/html",
                   "Dibuka di peramban mana pun; dapat dicetak menjadi PDF."),
    "md": Format("md", "Markdown", "md", "text/markdown",
                 "Teks polos untuk disunting lebih lanjut."),
    "json": Format("json", "JSON", "json", "application/json",
                   "Data terstruktur untuk diolah sistem lain."),
    "py": Format("py", "Sintaks Python", "py", "text/x-python",
                 "Skrip yang menjalankan ulang analisis dengan pustaka yang sama."),
    "r": Format("r", "Sintaks R", "R", "text/plain",
                "Skrip R untuk memeriksa silang hasil di luar aplikasi."),
    "zip": Format("zip", "Paket lengkap (ZIP)", "zip", "application/zip",
                  "Seluruh format, tabel dalam CSV, dan sintaks Python serta R."),
}

_PEMBUAT = {
    "docx": ke_docx,
    "pdf": ke_pdf,
    "xlsx": ke_xlsx,
    "pptx": ke_pptx,
    "html": ke_html,
    "md": ke_markdown,
    "json": ke_json,
    "py": ke_python,
    "r": ke_r,
    "zip": ke_zip,
}


def bangun(
    laporan: Laporan, kode_format: str, pembaca: str = "eksekutif", lengkap: bool = False
) -> bytes:
    """Hasilkan berkas laporan dalam format yang diminta."""
    if kode_format not in _PEMBUAT:
        raise ValueError(
            f"Format '{kode_format}' tidak dikenal. Pilih dari: {', '.join(FORMAT)}."
        )
    if pembaca not in AUDIENCES:
        raise ValueError(f"Pembaca '{pembaca}' tidak dikenal.")
    return _PEMBUAT[kode_format](laporan, pembaca, lengkap)


def nama_berkas(
    laporan: Laporan, kode_format: str, pembaca: str = "eksekutif", lengkap: bool = False
) -> str:
    if kode_format in {"py", "r"}:
        ragam = "sintaks_analisis"
    else:
        ragam = "laporan_lengkap" if lengkap else f"ringkasan_{pembaca}"
    dasar = str(laporan.dataset).rsplit(".", 1)[0]
    bersih = "".join(c if c.isalnum() or c in "-_" else "_" for c in dasar)[:40].strip("_")
    return f"{bersih or 'lentera'}_{ragam}.{FORMAT[kode_format].ekstensi}"
