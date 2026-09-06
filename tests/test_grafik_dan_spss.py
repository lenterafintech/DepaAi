"""Uji grafik statis untuk laporan dan pemuatan berkas SPSS."""

from __future__ import annotations

import io
import tempfile
from pathlib import Path

import pandas as pd
import pytest

from nalardata import ekspor, grafik, io_utils
from nalardata import narrative as nr

ROOT = Path(__file__).resolve().parents[1]
NUMERIK = [
    "usia",
    "pendapatan_bulanan",
    "rasio_utang_pendapatan",
    "skor_kredit",
    "jumlah_keterlambatan",
]


@pytest.fixture(scope="module")
def data() -> pd.DataFrame:
    return pd.read_csv(ROOT / "data" / "contoh_data_nasabah.csv")


@pytest.fixture(scope="module")
def laporan(data) -> nr.Laporan:
    konfig = nr.Konfigurasi(
        variabel=NUMERIK,
        nama_data="contoh_data_nasabah.csv",
        target_numerik="skor_kredit",
        prediktor=["rasio_utang_pendapatan", "jumlah_keterlambatan"],
        kelompok="segmen_usaha",
    )
    _, lap = nr.analisis_dan_laporan(data, konfig)
    return lap


def _png_sah(isi: bytes) -> bool:
    """Delapan byte pertama berkas PNG adalah tanda tangan tetap."""
    return isi.startswith(b"\x89PNG\r\n\x1a\n")


# --------------------------------------------------------------------------- #
# Grafik
# --------------------------------------------------------------------------- #


def test_peringkat_pendorong_menghasilkan_png(laporan):
    isi = grafik.peringkat_pendorong(laporan.pendorong)
    assert _png_sah(isi)
    assert len(isi) > 5_000


def test_status_pemeriksaan_menghasilkan_png(laporan):
    assert _png_sah(grafik.status_pemeriksaan(laporan.lampu))


def test_peta_korelasi_menghasilkan_png(data):
    assert _png_sah(grafik.peta_korelasi(data[NUMERIK].corr()))


def test_peta_korelasi_menerima_matriks_besar(data):
    """Angka dalam sel disembunyikan pada matriks besar agar tidak bertumpuk."""
    besar = data.select_dtypes(include="number")
    assert besar.shape[1] > 12
    assert _png_sah(grafik.peta_korelasi(besar.corr()))


@pytest.mark.parametrize(
    "fungsi, argumen, pesan",
    [
        (grafik.peringkat_pendorong, [], "Tidak ada pendorong"),
        (grafik.status_pemeriksaan, [], "Tidak ada pemeriksaan"),
        (grafik.peta_korelasi, pd.DataFrame(), "kosong"),
    ],
)
def test_masukan_kosong_ditolak(fungsi, argumen, pesan):
    with pytest.raises(ValueError, match=pesan):
        fungsi(argumen)


def test_grafik_laporan_mengembalikan_pasangan_judul_dan_gambar(laporan):
    hasil = grafik.grafik_laporan(laporan)
    assert hasil
    for judul, isi in hasil:
        assert isinstance(judul, str) and judul
        assert _png_sah(isi)


def test_kegagalan_grafik_tidak_menggagalkan_laporan():
    """Laporan tanpa isi tetap menghasilkan daftar kosong, bukan galat."""
    kosong = nr.Laporan(dataset="d.csv", n_baris=0, n_kolom=0, tanggal="01-01-2026")
    assert grafik.grafik_laporan(kosong) == []


# --------------------------------------------------------------------------- #
# Grafik pada berkas laporan
# --------------------------------------------------------------------------- #


def test_grafik_masuk_ke_word(laporan):
    from docx import Document

    dok = Document(io.BytesIO(ekspor.bangun(laporan, "docx", "akademik", True)))
    assert len(dok.inline_shapes) >= 2


def test_grafik_masuk_ke_pdf(laporan):
    from pypdf import PdfReader

    baca = PdfReader(io.BytesIO(ekspor.bangun(laporan, "pdf", "akademik", True)))
    assert sum(len(halaman.images) for halaman in baca.pages) >= 2


def test_grafik_masuk_ke_powerpoint(laporan):
    from pptx import Presentation

    presentasi = Presentation(io.BytesIO(ekspor.bangun(laporan, "pptx", "eksekutif")))
    gambar = sum(1 for s in presentasi.slides for bentuk in s.shapes if bentuk.shape_type == 13)
    assert gambar >= 2


def test_grafik_disematkan_ke_html(laporan):
    halaman = ekspor.bangun(laporan, "html", "akademik", True).decode("utf-8")
    # Gambar disematkan sebagai data URI agar berkasnya tetap mandiri.
    assert halaman.count("data:image/png;base64,") >= 2


def test_markdown_dan_json_menyebut_gambar_tanpa_menyematkannya(laporan):
    import json

    teks = ekspor.bangun(laporan, "md", "akademik", True).decode("utf-8")
    assert "grafik tersedia pada berkas" in teks

    isi = json.loads(ekspor.bangun(laporan, "json", "akademik", True))
    assert any(b["ada_gambar"] for b in isi["blok"])
    # JSON tetap ringan karena gambar tidak ikut disematkan.
    assert "base64" not in json.dumps(isi)


# --------------------------------------------------------------------------- #
# Berkas SPSS
# --------------------------------------------------------------------------- #


@pytest.fixture
def berkas_sav(tmp_path) -> Path:
    import pyreadstat

    df = pd.DataFrame(
        {"usia": [30, 45, 52], "jk": [1, 2, 1], "skor": [610.5, 720.0, 588.25]}
    )
    jalur = tmp_path / "uji.sav"
    pyreadstat.write_sav(
        df,
        str(jalur),
        variable_value_labels={"jk": {1: "Laki-laki", 2: "Perempuan"}},
        column_labels=["Usia responden", "Jenis kelamin", "Skor kredit"],
    )
    return jalur


def test_sav_terbaca_dari_jalur_berkas(berkas_sav):
    hasil = io_utils.load_table(berkas_sav)
    assert list(hasil.columns) == ["usia", "jk", "skor"]
    assert len(hasil) == 3


def test_label_nilai_spss_diterapkan(berkas_sav):
    """Tanpa label, kolom kategori muncul sebagai angka tanpa makna."""
    hasil = io_utils.load_table(berkas_sav)
    assert hasil["jk"].tolist() == ["Laki-laki", "Perempuan", "Laki-laki"]


def test_sav_terbaca_dari_aliran_unggahan(berkas_sav):
    """Unggahan Streamlit berupa aliran dalam memori, bukan jalur berkas."""
    aliran = io.BytesIO(berkas_sav.read_bytes())
    aliran.name = "uji.sav"
    hasil = io_utils.load_table(aliran, filename="uji.sav")
    assert hasil.equals(io_utils.load_table(berkas_sav))


def test_berkas_sementara_dibersihkan_setelah_membaca_aliran(berkas_sav):
    sebelum = set(Path(tempfile.gettempdir()).glob("*.sav"))
    aliran = io.BytesIO(berkas_sav.read_bytes())
    io_utils.load_table(aliran, filename="uji.sav")
    sesudah = set(Path(tempfile.gettempdir()).glob("*.sav"))
    assert sesudah == sebelum


def test_kamus_data_spss(berkas_sav):
    kamus = io_utils.meta_spss(berkas_sav)
    assert list(kamus["Kolom"]) == ["usia", "jk", "skor"]
    assert kamus.set_index("Kolom").loc["jk", "Label"] == "Jenis kelamin"


def test_sav_terdaftar_sebagai_format_yang_didukung():
    assert ".sav" in io_utils.SUPPORTED_SUFFIXES
    assert ".zsav" in io_utils.SUPPORTED_SUFFIXES


def test_format_asing_tetap_ditolak(tmp_path):
    berkas = tmp_path / "data.dta"
    berkas.write_bytes(b"x")
    with pytest.raises(io_utils.UnsupportedFileError, match="belum didukung"):
        io_utils.load_table(berkas)
