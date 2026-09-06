"""Ekspor laporan hasil analisis ke berbagai format berkas.

Laporan disusun sekali menjadi rangkaian blok (judul, paragraf, poin, tabel), lalu
tiap format hanya menerjemahkan blok yang sama. Dengan begitu isi Word, PDF,
PowerPoint, Excel, HTML, dan Markdown tidak pernah berbeda satu sama lain.

Dua ragam laporan tersedia:

- **Ringkasan** — satu register pembaca saja, ringkas untuk dibagikan.
- **Laporan lengkap** — ketiga register, seluruh tabel, paragraf siap salin,
  rujukan, dan catatan metode.
"""

from __future__ import annotations

import io
import json
import zipfile
from html import escape
from datetime import date
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd

from mv_statlab.keranjang import Keranjang
from mv_statlab.narrative import AUDIENCE_LABELS, AUDIENCES, Laporan, tabel_markdown
from mv_statlab.report_html import laporan_html, laporan_html_semua
from mv_statlab.sintaks import bangkitkan

# Padanan aman untuk lambang yang tidak tersedia pada huruf bawaan PDF.
GANTI_LAMBANG = {
    "−": "-",
    "–": "-",
    "—": "-",
    "×": "x",
    "√": "akar ",
    "≥": ">=",
    "≤": "<=",
    "η": "eta",
    "χ": "chi",
    "β": "beta",
    "λ": "lambda",
    "α": "alpha",
    "²": "^2",
    "·": "-",
    "“": '"',
    "”": '"',
    "’": "'",
}
JALUR_FONT = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/DejaVuSans.ttf",
    "C:/Windows/Fonts/DejaVuSans.ttf",
]


@dataclass
class Blok:
    """Satu potongan isi laporan, bebas dari format berkasnya."""

    jenis: str  # judul | subjudul | paragraf | poin | tabel | gambar | meta | catatan
    teks: str = ""
    poin: list[str] = field(default_factory=list)
    tabel: pd.DataFrame | None = None
    gambar: bytes | None = None
    catatan: str = ""


@dataclass
class Dokumen:
    """Isi laporan yang sudah lepas dari sumbernya.

    Lapisan ini yang memungkinkan satu berkas ekspor dibangun baik dari laporan
    naratif maupun dari keranjang hasil, tanpa tiap penulis format perlu tahu
    asal-usulnya.
    """

    judul: str
    meta: str = ""
    blok: list[Blok] = field(default_factory=list)
    nama_dasar: str = "laporan"

    def tabel(self) -> list[tuple[str, pd.DataFrame]]:
        """Pasangan (judul, tabel) menurut subjudul terdekat sebelumnya."""
        hasil: list[tuple[str, pd.DataFrame]] = []
        tajuk = self.judul
        for b in self.blok:
            if b.jenis in ("subjudul", "paragraf") and b.teks:
                tajuk = b.teks
            elif b.jenis == "tabel" and b.tabel is not None:
                hasil.append((tajuk, b.tabel))
        return hasil


# --------------------------------------------------------------------------- #
# Penyusun isi
# --------------------------------------------------------------------------- #


def susun_blok(laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False) -> list[Blok]:
    """Ubah laporan menjadi rangkaian blok yang siap ditulis ke format apa pun.

    Dua sumbu yang saling bebas menentukan isinya:

    - ``pembaca`` menentukan **bahasanya** — eksekutif, akademik, atau profesional.
      Berapa pun kedalamannya, laporan ditulis dalam satu register saja; menggabung
      ketiganya justru membuat pembaca membaca hal yang sama tiga kali.
    - ``lengkap`` menentukan **kedalamannya**. Ringkasan memuat kesimpulan, status
      pemeriksaan, pendorong, dan rekomendasi. Laporan lengkap menambahkan uraian
      tiap temuan, seluruh tabel hasil, kalimat siap salin, rujukan ambang, dan
      catatan analisis yang tidak dapat dijalankan.
    """
    if pembaca not in AUDIENCES:
        raise ValueError(f"Pembaca '{pembaca}' tidak dikenal. Pilih dari {AUDIENCES}.")

    judul = (
        f"Laporan Lengkap Analisis Multivariat — {AUDIENCE_LABELS[pembaca]}"
        if lengkap
        else f"Ringkasan {AUDIENCE_LABELS[pembaca]}"
    )
    blok: list[Blok] = [
        Blok("judul", judul),
        Blok(
            "meta",
            f"Sumber data: {laporan.dataset} · {laporan.n_baris} baris × "
            f"{laporan.n_kolom} kolom · {len(laporan.metode_terpakai)} metode · "
            f"{laporan.tanggal}",
        ),
        Blok("paragraf", laporan.headline),
        Blok("paragraf", laporan.subheadline),
    ]

    gambar = {nama: isi for nama, isi in _grafik(laporan)}

    if laporan.lampu:
        blok.append(Blok("subjudul", "Status pemeriksaan"))
        blok.append(
            Blok(
                "tabel",
                tabel=pd.DataFrame(
                    {
                        "Pemeriksaan": [l.label for l in laporan.lampu],
                        "Nilai": [l.nilai for l in laporan.lampu],
                        "Status": [l.status_label for l in laporan.lampu],
                        "Catatan": [l.catatan for l in laporan.lampu],
                    }
                ),
            )
        )
        if "Status pemeriksaan" in gambar:
            blok.append(Blok("gambar", gambar=gambar["Status pemeriksaan"]))

    if laporan.pendorong:
        blok.append(Blok("subjudul", "Peringkat pendorong"))
        blok.append(Blok("catatan", catatan=f"Sumber: {laporan.pendorong_sumber}"))
        blok.append(
            Blok(
                "tabel",
                tabel=pd.DataFrame(
                    {
                        "Faktor": [p.nama for p in laporan.pendorong],
                        "Satuan": [p.satuan for p in laporan.pendorong],
                        "Nilai": [round(p.nilai, 4) for p in laporan.pendorong],
                        "Arah": [p.arah for p in laporan.pendorong],
                        "Signifikan": ["Ya" if p.signifikan else "Tidak" for p in laporan.pendorong],
                        "Catatan": [p.catatan for p in laporan.pendorong],
                    }
                ),
            )
        )
        if "Peringkat pendorong" in gambar:
            blok.append(Blok("gambar", gambar=gambar["Peringkat pendorong"]))

    blok.append(Blok("subjudul", "Temuan"))
    if lengkap:
        # Laporan lengkap menguraikan tiap temuan beserta metodenya.
        for temuan in laporan.temuan:
            blok.append(Blok("paragraf", f"{temuan.judul} ({temuan.metode})"))
            blok.append(Blok("paragraf", temuan.teks(pembaca)))
    else:
        # Ringkasan cukup memuat inti tiap temuan dalam satu daftar.
        blok.append(Blok("poin", poin=[t.ringkas for t in laporan.temuan]))

    if lengkap:
        for nomor, (judul_tabel, tabel, catatan) in laporan.tabel.items():
            blok.append(Blok("subjudul", f"{nomor}. {judul_tabel}"))
            blok.append(Blok("tabel", tabel=tabel, catatan=catatan))

        if laporan.paragraf and pembaca == "akademik":
            # Kalimat siap salin mengikuti konvensi pelaporan statistik, sehingga
            # hanya bermakna pada register akademik.
            blok.append(Blok("subjudul", "Kalimat siap salin"))
            for paragraf in laporan.paragraf:
                blok.append(Blok("paragraf", paragraf.bagian))
                blok.append(Blok("paragraf", paragraf.teks))

    if laporan.rekomendasi:
        blok.append(Blok("subjudul", "Rekomendasi tindakan"))
        blok.append(
            Blok(
                "poin",
                poin=[
                    f"{r.judul} (prioritas {r.prioritas}) — {r.alasan}"
                    for r in laporan.rekomendasi
                ],
            )
        )

    if laporan.keterbatasan:
        blok.append(Blok("subjudul", "Batas kesimpulan"))
        blok.append(Blok("poin", poin=list(laporan.keterbatasan)))

    if lengkap and laporan.rujukan:
        blok.append(Blok("subjudul", "Rujukan ambang yang dipakai"))
        blok.append(Blok("poin", poin=list(laporan.rujukan)))

    if lengkap and laporan.dilewati:
        blok.append(Blok("subjudul", "Analisis yang tidak dapat dijalankan"))
        blok.append(Blok("poin", poin=list(laporan.dilewati)))

    blok.append(
        Blok(
            "catatan",
            catatan=(
                "Disusun otomatis oleh MV Statlab. Kesimpulan statistik menunjukkan pola "
                "dalam data, bukan bukti sebab-akibat; keputusan akhir tetap memerlukan "
                "pertimbangan konteks dan keahlian bidang."
            ),
        )
    )
    return blok


def dari_laporan(
    laporan: Laporan, pembaca: str = "eksekutif", lengkap: bool = False
) -> Dokumen:
    """Dokumen dari laporan naratif, memakai penyusun blok yang sudah ada."""
    blok = susun_blok(laporan, pembaca, lengkap)
    judul = next((b.teks for b in blok if b.jenis == "judul"), "Laporan")
    meta = next((b.teks for b in blok if b.jenis == "meta"), "")
    dasar = "laporan_lengkap" if lengkap else f"ringkasan_{pembaca}"
    return Dokumen(judul=judul, meta=meta, blok=blok, nama_dasar=dasar)


def dari_keranjang(keranjang: Keranjang) -> Dokumen:
    """Dokumen dari hasil yang dikumpulkan pengguna di halaman-halaman metode."""
    if keranjang.kosong():
        raise ValueError(
            "Keranjang hasil masih kosong. Jalankan analisis lebih dulu, lalu tekan "
            "'Simpan ke laporan' pada hasil yang ingin dilaporkan."
        )

    ringkas = keranjang.ringkas()
    meta = (
        f"{ringkas['Bagian']} bagian · {ringkas['Tabel']} tabel · "
        f"{ringkas['Tafsiran']} tafsiran · disusun {date.today().strftime('%d-%m-%Y')}"
    )
    if keranjang.peneliti.strip():
        meta = f"Disusun oleh {keranjang.peneliti.strip()} · {meta}"

    blok: list[Blok] = [
        Blok("judul", keranjang.judul),
        Blok("meta", meta),
        Blok("subjudul", "Daftar isi"),
        Blok("tabel", tabel=keranjang.daftar_isi()),
    ]

    for nama, isi in keranjang.per_bagian().items():
        blok.append(Blok("subjudul", nama))
        for item in isi:
            if item.jenis == "tabel" and item.tabel is not None:
                blok.append(Blok("paragraf", item.judul))
                blok.append(Blok("tabel", tabel=item.tabel, catatan=item.catatan))
            else:
                blok.append(Blok("paragraf", item.teks))

    blok.append(
        Blok(
            "catatan",
            catatan=(
                "Laporan ini memuat hasil analisis yang Anda jalankan dan simpan "
                "sendiri. Kesimpulan statistik menunjukkan pola dalam data, bukan "
                "bukti sebab-akibat."
            ),
        )
    )
    return Dokumen(judul=keranjang.judul, meta=meta, blok=blok, nama_dasar="laporan_hasil")


def _dokumen(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> Dokumen:
    """Terima Laporan maupun Keranjang, kembalikan dokumen yang siap ditulis."""
    if isinstance(sumber, Dokumen):
        return sumber
    if isinstance(sumber, Keranjang):
        return dari_keranjang(sumber)
    return dari_laporan(sumber, pembaca, lengkap)


def _grafik(laporan: Laporan) -> list[tuple[str, bytes]]:
    """Grafik statis untuk laporan; kegagalannya tidak menggagalkan ekspor."""
    try:
        from mv_statlab.grafik import grafik_laporan

        return grafik_laporan(laporan)
    except Exception:  # noqa: BLE001 - matplotlib tak ada atau render gagal
        return []


def _aman_pdf(teks: str) -> str:
    """Ganti lambang yang tidak ada pada huruf bawaan PDF dengan padanan ASCII."""
    hasil = str(teks)
    for lambang, ganti in GANTI_LAMBANG.items():
        hasil = hasil.replace(lambang, ganti)
    return hasil


# --------------------------------------------------------------------------- #
# Word (.docx)
# --------------------------------------------------------------------------- #


def ke_docx(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    dokumen = _dokumen(sumber, pembaca, lengkap)
    dok = Document()
    gaya = dok.styles["Normal"]
    gaya.font.name = "Calibri"
    gaya.font.size = Pt(11)

    for blok in dokumen.blok:
        if blok.jenis == "judul":
            dok.add_heading(blok.teks, level=0)
        elif blok.jenis == "subjudul":
            dok.add_heading(blok.teks, level=1)
        elif blok.jenis == "meta":
            paragraf = dok.add_paragraph(blok.teks)
            paragraf.runs[0].font.size = Pt(9)
            paragraf.runs[0].font.color.rgb = RGBColor(0x6F, 0x7A, 0x91)
        elif blok.jenis == "paragraf":
            dok.add_paragraph(blok.teks)
        elif blok.jenis == "poin":
            for butir in blok.poin:
                dok.add_paragraph(butir, style="List Bullet")
        elif blok.jenis == "catatan" and blok.catatan:
            paragraf = dok.add_paragraph(blok.catatan)
            paragraf.runs[0].italic = True
            paragraf.runs[0].font.size = Pt(9)
        elif blok.jenis == "gambar" and blok.gambar:
            from docx.shared import Inches

            dok.add_picture(io.BytesIO(blok.gambar), width=Inches(6.2))
        elif blok.jenis == "tabel" and blok.tabel is not None:
            tabel = blok.tabel.fillna("")
            objek = dok.add_table(rows=1, cols=len(tabel.columns))
            objek.style = "Light Grid Accent 1"
            for i, kolom in enumerate(tabel.columns):
                sel = objek.rows[0].cells[i]
                sel.text = str(kolom)
                for p in sel.paragraphs:
                    for r in p.runs:
                        r.bold = True
            for _, baris in tabel.iterrows():
                sel_baris = objek.add_row().cells
                for i, nilai in enumerate(baris):
                    sel_baris[i].text = str(nilai)
            if blok.catatan:
                catatan = dok.add_paragraph(f"Catatan. {blok.catatan}")
                catatan.runs[0].italic = True
                catatan.runs[0].font.size = Pt(9)

    penampung = io.BytesIO()
    dok.save(penampung)
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# Excel (.xlsx)
# --------------------------------------------------------------------------- #


def _nama_lembar(teks: str, dipakai: set[str]) -> str:
    """Nama lembar Excel: maksimal 31 karakter, tanpa lambang terlarang, dan unik."""
    bersih = "".join(c for c in str(teks) if c not in r"[]:*?/\\")[:31].strip() or "Lembar"
    calon = bersih
    urut = 2
    while calon.lower() in dipakai:
        akhiran = f" ({urut})"
        calon = bersih[: 31 - len(akhiran)] + akhiran
        urut += 1
    dipakai.add(calon.lower())
    return calon


def ke_xlsx(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    return _dok_ke_xlsx(_dokumen(sumber, pembaca, lengkap))


# --------------------------------------------------------------------------- #
# PDF
# --------------------------------------------------------------------------- #


def _daftarkan_font() -> str:
    """Pakai DejaVu bila tersedia agar lambang statistik tampil utuh."""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for jalur in JALUR_FONT:
        berkas = Path(jalur)
        if berkas.exists():
            try:
                pdfmetrics.registerFont(TTFont("DejaVu", str(berkas)))
                tebal = berkas.with_name(berkas.stem + "-Bold" + berkas.suffix)
                if tebal.exists():
                    pdfmetrics.registerFont(TTFont("DejaVu-Bold", str(tebal)))
                    pdfmetrics.registerFontFamily("DejaVu", normal="DejaVu", bold="DejaVu-Bold")
                return "DejaVu"
            except Exception:  # noqa: BLE001 - huruf gagal dimuat, pakai bawaan
                break
    return "Helvetica"


def ke_pdf(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    # Nama 'dokumen' sudah dipakai SimpleDocTemplate di bawah, jadi isi laporan
    # disimpan dengan nama lain agar tidak tertimpa.
    konten = _dokumen(sumber, pembaca, lengkap)
    font = _daftarkan_font()
    font_tebal = "DejaVu-Bold" if font == "DejaVu" else "Helvetica-Bold"
    perlu_sanitasi = font == "Helvetica"

    def teks(nilai: object) -> str:
        mentah = str(nilai)
        return _aman_pdf(mentah) if perlu_sanitasi else mentah

    dasar = getSampleStyleSheet()
    gaya_judul = ParagraphStyle(
        "JudulStatlab", parent=dasar["Title"], fontName=font_tebal, fontSize=18,
        textColor=colors.HexColor("#131a2b"), alignment=TA_LEFT, spaceAfter=6,
    )
    gaya_sub = ParagraphStyle(
        "SubStatlab", parent=dasar["Heading2"], fontName=font_tebal, fontSize=13,
        textColor=colors.HexColor("#26356b"), spaceBefore=14, spaceAfter=6,
    )
    gaya_isi = ParagraphStyle(
        "IsiStatlab", parent=dasar["BodyText"], fontName=font, fontSize=10, leading=15,
        textColor=colors.HexColor("#1f2937"), spaceAfter=6,
    )
    gaya_meta = ParagraphStyle(
        "MetaStatlab", parent=gaya_isi, fontSize=8.5,
        textColor=colors.HexColor("#6f7a91"),
    )

    penampung = io.BytesIO()
    dokumen = SimpleDocTemplate(
        penampung, pagesize=A4,
        leftMargin=20 * mm, rightMargin=18 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
        title=konten.judul,
        author="MV Statlab",
    )
    isi: list = []
    for blok in konten.blok:
        if blok.jenis == "judul":
            isi.append(Paragraph(teks(blok.teks), gaya_judul))
        elif blok.jenis == "subjudul":
            isi.append(Paragraph(teks(blok.teks), gaya_sub))
        elif blok.jenis == "meta":
            isi.append(Paragraph(teks(blok.teks), gaya_meta))
        elif blok.jenis == "paragraf":
            isi.append(Paragraph(teks(blok.teks), gaya_isi))
        elif blok.jenis == "poin":
            for butir in blok.poin:
                isi.append(Paragraph("• " + teks(butir), gaya_isi))
        elif blok.jenis == "catatan" and blok.catatan:
            isi.append(Paragraph(teks(blok.catatan), gaya_meta))
        elif blok.jenis == "gambar" and blok.gambar:
            from reportlab.platypus import Image as GambarPDF

            try:
                gambar_pdf = GambarPDF(io.BytesIO(blok.gambar))
                skala = dokumen.width / gambar_pdf.imageWidth
                gambar_pdf.drawWidth = dokumen.width
                gambar_pdf.drawHeight = gambar_pdf.imageHeight * skala
                isi.append(gambar_pdf)
                isi.append(Spacer(1, 10))
            except Exception:  # noqa: BLE001 - gambar rusak dilewati
                pass
        elif blok.jenis == "tabel" and blok.tabel is not None:
            tabel = blok.tabel.fillna("")
            data = [[Paragraph(f"<b>{teks(k)}</b>", gaya_meta) for k in tabel.columns]]
            for _, baris in tabel.iterrows():
                data.append([Paragraph(teks(v), gaya_meta) for v in baris])
            lebar = (dokumen.width) / max(len(tabel.columns), 1)
            objek = Table(data, colWidths=[lebar] * len(tabel.columns), repeatRows=1)
            objek.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eef1f8")),
                        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#dde3ee")),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 4),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                        ("TOPPADDING", (0, 0), (-1, -1), 3),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                    ]
                )
            )
            isi.append(objek)
            if blok.catatan:
                isi.append(Spacer(1, 3))
                isi.append(Paragraph("Catatan. " + teks(blok.catatan), gaya_meta))
            isi.append(Spacer(1, 8))

    dokumen.build(isi)
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# PowerPoint (.pptx)
# --------------------------------------------------------------------------- #


def ke_pptx(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    return _dok_ke_pptx(_dokumen(sumber, pembaca, lengkap))


# --------------------------------------------------------------------------- #
# Format berbasis teks dan paket data
# --------------------------------------------------------------------------- #


def ke_html(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    return _dok_ke_html(_dokumen(sumber, pembaca, lengkap))


def ke_markdown(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    return _dok_ke_markdown(_dokumen(sumber, pembaca, lengkap))


def ke_json(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Isi laporan dalam bentuk terstruktur, untuk diolah kembali oleh sistem lain."""
    return _dok_ke_json(_dokumen(sumber, pembaca, lengkap))


def ke_python(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Skrip Python yang menjalankan ulang analisis dengan pustaka yang sama."""
    return bangkitkan(getattr(sumber, "konfig", None), "py").encode("utf-8")


def ke_r(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Skrip R sebagai pemeriksaan silang di luar Python."""
    return bangkitkan(getattr(sumber, "konfig", None), "r").encode("utf-8")


def _petunjuk_sintaks() -> str:
    """Keterangan yang menyertai berkas sintaks di dalam paket ZIP."""
    return (
        "SINTAKS UNTUK PERANGKAT LAIN\n"
        "============================\n\n"
        "Berkas di folder ini menjalankan ulang analisis yang sama di perangkat lain:\n\n"
        "  analisis.py     Python - memakai pustaka yang sama, angkanya identik\n"
        "  analisis.R      R - padanan terdekat, untuk pemeriksaan silang\n"
        "  analisis.sps    SPSS - perintah yang setara\n"
        "  model_amos.txt  AMOS - spesifikasi yang perlu digambar ulang\n"
        "  analisis.inp    Mplus - berkas input siap jalan\n\n"
        "SEBELUM MENJALANKAN: seluruh berkas merujuk data.csv, yang TIDAK disertakan\n"
        "di sini karena paket laporan memang tidak membawa data mentah. Ambil datanya\n"
        "dengan salah satu cara berikut, lalu letakkan sebagai data.csv di folder yang\n"
        "sama:\n\n"
        "  - Unduh dari halaman Beranda & Data pada aplikasi, atau\n"
        "  - Buka berkas proyek .mvstatlab Anda; data.csv ada di dalamnya.\n\n"
        "Perbedaan kecil antar perangkat adalah hal wajar. Sebelum membandingkan,\n"
        "samakan lebih dulu penanganan nilai hilang, standardisasi, pengkodean\n"
        "kategori, dan tipe jumlah kuadrat.\n"
    )


def ke_spss(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Sintaks SPSS yang setara, untuk memeriksa ulang hasil di perangkat itu."""
    return bangkitkan(getattr(sumber, "konfig", None), "spss").encode("utf-8")


def ke_amos(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Spesifikasi model untuk dipindahkan ke AMOS."""
    return bangkitkan(getattr(sumber, "konfig", None), "amos").encode("utf-8")


def ke_mplus(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Berkas input Mplus (.inp)."""
    return bangkitkan(getattr(sumber, "konfig", None), "mplus").encode("utf-8")


def ke_zip(sumber, pembaca: str = "eksekutif", lengkap: bool = False) -> bytes:
    """Paket lengkap: laporan dalam beberapa format sekaligus, ditambah tabel CSV.

    Sintaks Python dan R hanya dapat dibangkitkan bila sumbernya membawa konfigurasi
    analisis; keranjang hasil tidak membawanya, sehingga paketnya berisi berkas
    laporan dan tabel saja.
    """
    dokumen = _dokumen(sumber, pembaca, lengkap)
    konfig = getattr(sumber, "konfig", None)
    return _dok_ke_zip(dokumen, konfig)


# --------------------------------------------------------------------------- #
# Penulis berbasis blok
# --------------------------------------------------------------------------- #
#
# Jalur laporan naratif tetap memakai penulis khususnya, karena kekayaan bentuknya
# — panel tiga register pada HTML, lembar per bagian pada Excel — tidak terwakili
# oleh blok. Penulis di bawah melayani dokumen mana pun, termasuk keranjang hasil.


def _dok_ke_xlsx(dokumen: Dokumen) -> bytes:
    """Tiap tabel menjadi lembar tersendiri, ditambah lembar keterangan."""
    penampung = io.BytesIO()
    dipakai: set[str] = set()
    with pd.ExcelWriter(penampung, engine="openpyxl") as penulis:
        pd.DataFrame(
            {
                "Keterangan": ["Judul", "Keterangan", "Jumlah tabel"],
                "Isi": [
                    dokumen.judul,
                    dokumen.meta,
                    str(len(dokumen.tabel())),
                ],
            }
        ).to_excel(penulis, sheet_name=_nama_lembar("Keterangan", dipakai), index=False)

        for judul, tabel in dokumen.tabel():
            tabel.to_excel(
                penulis, sheet_name=_nama_lembar(judul, dipakai), index=False
            )

        teks = [b.teks for b in dokumen.blok if b.jenis == "paragraf" and b.teks]
        if teks:
            pd.DataFrame({"Catatan": teks}).to_excel(
                penulis, sheet_name=_nama_lembar("Catatan", dipakai), index=False
            )
    return penampung.getvalue()


def _dok_ke_pptx(dokumen: Dokumen) -> bytes:
    """Satu slide judul, lalu satu slide per bagian."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentasi = Presentation()
    presentasi.slide_width = Inches(13.333)
    presentasi.slide_height = Inches(7.5)
    navy = RGBColor(0x26, 0x35, 0x6B)

    slide = presentasi.slides.add_slide(presentasi.slide_layouts[0])
    slide.shapes.title.text = dokumen.judul
    slide.placeholders[1].text = dokumen.meta

    def slide_isi(judul: str, butir: list[str], maks: int = 6) -> None:
        for mulai in range(0, len(butir), maks):
            s = presentasi.slides.add_slide(presentasi.slide_layouts[1])
            s.shapes.title.text = judul if mulai == 0 else f"{judul} (lanjutan)"
            if s.shapes.title.text_frame.paragraphs[0].runs:
                s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
            kerangka = s.placeholders[1].text_frame
            kerangka.clear()
            kerangka.word_wrap = True
            for i, isi in enumerate(butir[mulai : mulai + maks]):
                paragraf = kerangka.paragraphs[0] if i == 0 else kerangka.add_paragraph()
                paragraf.text = isi[:400]
                paragraf.font.size = Pt(15)

    # Blok dikelompokkan menurut subjudul terdekat agar tiap bagian menjadi satu slide.
    bagian: dict[str, list[str]] = {}
    slide_gambar: list[tuple[str, bytes]] = []
    tajuk = dokumen.judul
    for b in dokumen.blok:
        if b.jenis == "subjudul" and b.teks:
            tajuk = b.teks
            bagian.setdefault(tajuk, [])
        elif b.jenis == "paragraf" and b.teks:
            bagian.setdefault(tajuk, []).append(b.teks)
        elif b.jenis == "poin" and b.poin:
            bagian.setdefault(tajuk, []).extend(b.poin)
        elif b.jenis == "gambar" and b.gambar:
            slide_gambar.append((tajuk, b.gambar))
        elif b.jenis == "tabel" and b.tabel is not None:
            bagian.setdefault(tajuk, []).append(
                f"[tabel {b.tabel.shape[0]} baris x {b.tabel.shape[1]} kolom - "
                "lihat berkas Word, Excel, atau PDF]"
            )

    for judul, butir in bagian.items():
        if butir:
            slide_isi(judul, butir)

    # Grafik mendapat slide sendiri agar terbaca pada layar presentasi.
    for judul, isi_gambar in slide_gambar:
        s = presentasi.slides.add_slide(presentasi.slide_layouts[5])
        s.shapes.title.text = judul
        try:
            s.shapes.add_picture(
                io.BytesIO(isi_gambar), Inches(0.6), Inches(1.5), width=Inches(12.1)
            )
        except Exception:  # noqa: BLE001 - gambar rusak dilewati
            pass

    penampung = io.BytesIO()
    presentasi.save(penampung)
    return penampung.getvalue()


def _dok_ke_html(dokumen: Dokumen) -> bytes:
    """Halaman mandiri sederhana yang mengikuti tema terang/gelap peramban."""
    bagian: list[str] = []
    for b in dokumen.blok:
        if b.jenis == "judul":
            bagian.append(f"<h1>{escape(b.teks)}</h1>")
        elif b.jenis == "meta":
            bagian.append(f'<p class="meta">{escape(b.teks)}</p>')
        elif b.jenis == "subjudul":
            bagian.append(f"<h2>{escape(b.teks)}</h2>")
        elif b.jenis == "paragraf" and b.teks:
            bagian.append(f"<p>{escape(b.teks)}</p>")
        elif b.jenis == "poin" and b.poin:
            butir = "".join(f"<li>{escape(str(p))}</li>" for p in b.poin)
            bagian.append(f"<ul>{butir}</ul>")
        elif b.jenis == "catatan" and b.catatan:
            bagian.append(f'<p class="catatan">{escape(b.catatan)}</p>')
        elif b.jenis == "gambar" and b.gambar:
            import base64

            sandi = base64.b64encode(b.gambar).decode("ascii")
            bagian.append(
                f'<img alt="Grafik laporan" src="data:image/png;base64,{sandi}">'
            )
        elif b.jenis == "tabel" and b.tabel is not None:
            tabel = b.tabel.fillna("")
            kepala = "".join(f"<th>{escape(str(k))}</th>" for k in tabel.columns)
            baris = "".join(
                "<tr>" + "".join(f"<td>{escape(str(v))}</td>" for v in nilai) + "</tr>"
                for nilai in tabel.itertuples(index=False)
            )
            bagian.append(
                f'<div class="tabel"><table><thead><tr>{kepala}</tr></thead>'
                f"<tbody>{baris}</tbody></table></div>"
            )
            if b.catatan:
                bagian.append(f'<p class="catatan">Catatan. {escape(b.catatan)}</p>')

    halaman = f"""<!doctype html>
<html lang="id"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(dokumen.judul)}</title>
<style>
:root{{color-scheme:light dark;--kertas:#fff;--tinta:#131a2b;--tinta2:#3d4860;
--redup:#6f7a91;--garis:#dde3ee;--aksen:#26356b;--samar:#f4f6fb}}
@media (prefers-color-scheme:dark){{:root:not([data-theme="light"]){{
--kertas:#0f1420;--tinta:#e9edf6;--tinta2:#b3bccf;--redup:#8b96ac;
--garis:#2a344c;--aksen:#93a6ea;--samar:#161d2e}}}}
*{{box-sizing:border-box}}
body{{margin:0;padding:32px 20px 64px;background:var(--kertas);color:var(--tinta);
font:15px/1.65 -apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif}}
main{{max-width:76ch;margin:0 auto}}
h1{{font-size:1.8rem;line-height:1.25;margin:0 0 .4rem;letter-spacing:-.01em}}
h2{{font-size:1.05rem;margin:2rem 0 .6rem;color:var(--aksen);
border-bottom:1px solid var(--garis);padding-bottom:.3rem}}
p{{margin:0 0 .9rem;color:var(--tinta2)}}
.meta{{font-family:ui-monospace,Menlo,monospace;font-size:.78rem;color:var(--redup);
margin-bottom:1.6rem}}
.catatan{{font-size:.84rem;color:var(--redup);font-style:italic}}
ul{{margin:0 0 1rem;padding-left:1.3rem;color:var(--tinta2)}}
li{{margin-bottom:.4rem}}
.tabel{{overflow-x:auto;border:1px solid var(--garis);border-radius:8px;
margin-bottom:.6rem}}
img{{max-width:100%;height:auto;display:block;margin:0 0 1rem;border:1px solid var(--garis);
border-radius:8px;background:#fff}}
table{{border-collapse:collapse;width:100%;font-size:.84rem}}
th{{background:var(--samar);text-align:left;padding:8px 11px;font-size:.72rem;
text-transform:uppercase;letter-spacing:.04em;color:var(--redup);
border-bottom:1px solid var(--garis);white-space:nowrap}}
td{{padding:7px 11px;border-bottom:1px solid var(--garis);color:var(--tinta2);
white-space:nowrap}}
tbody tr:last-child td{{border-bottom:0}}
</style></head><body><main>
{"".join(bagian)}
</main></body></html>"""
    return halaman.encode("utf-8")


def _dok_ke_markdown(dokumen: Dokumen) -> bytes:
    baris: list[str] = []
    for b in dokumen.blok:
        if b.jenis == "judul":
            baris += [f"# {b.teks}", ""]
        elif b.jenis == "meta":
            baris += [f"*{b.teks}*", ""]
        elif b.jenis == "subjudul":
            baris += [f"## {b.teks}", ""]
        elif b.jenis == "paragraf" and b.teks:
            baris += [b.teks, ""]
        elif b.jenis == "poin" and b.poin:
            baris += [f"- {p}" for p in b.poin] + [""]
        elif b.jenis == "catatan" and b.catatan:
            baris += [f"> {b.catatan}", ""]
        elif b.jenis == "gambar" and b.gambar:
            baris += ["*(grafik tersedia pada berkas Word, PDF, PowerPoint, dan HTML)*", ""]
        elif b.jenis == "tabel" and b.tabel is not None:
            baris += [tabel_markdown(b.tabel), ""]
            if b.catatan:
                baris += [f"*Catatan. {b.catatan}*", ""]
    return "\\n".join(baris).encode("utf-8")


def _dok_ke_json(dokumen: Dokumen) -> bytes:
    isi = {
        "judul": dokumen.judul,
        "meta": dokumen.meta,
        "blok": [
            {
                "jenis": b.jenis,
                "teks": b.teks,
                "poin": b.poin,
                "catatan": b.catatan,
                "tabel": (
                    None if b.tabel is None else b.tabel.fillna("").to_dict(orient="records")
                ),
                # Gambar tidak disematkan ke JSON agar berkasnya tetap ringan dan
                # dapat dibaca; keberadaannya cukup ditandai.
                "ada_gambar": bool(b.gambar),
            }
            for b in dokumen.blok
        ],
    }
    return json.dumps(isi, ensure_ascii=False, indent=2).encode("utf-8")


def _dok_ke_zip(dokumen: Dokumen, konfig=None) -> bytes:
    penampung = io.BytesIO()
    dasar = dokumen.nama_dasar
    with zipfile.ZipFile(penampung, "w", zipfile.ZIP_DEFLATED) as arsip:
        arsip.writestr(f"{dasar}.html", _dok_ke_html(dokumen))
        arsip.writestr(f"{dasar}.md", _dok_ke_markdown(dokumen))
        arsip.writestr(f"{dasar}.docx", ke_docx(dokumen))
        arsip.writestr(f"{dasar}.xlsx", _dok_ke_xlsx(dokumen))
        arsip.writestr(f"{dasar}.pdf", ke_pdf(dokumen))
        arsip.writestr(f"{dasar}.json", _dok_ke_json(dokumen))
        if konfig is not None:
            for berkas, bahasa in (
                ("sintaks/analisis.py", "py"),
                ("sintaks/analisis.R", "r"),
                ("sintaks/analisis.sps", "spss"),
                ("sintaks/model_amos.txt", "amos"),
                ("sintaks/analisis.inp", "mplus"),
            ):
                arsip.writestr(berkas, bangkitkan(konfig, bahasa))
            # Seluruh sintaks merujuk data.csv, yang tidak ikut di sini karena paket
            # laporan memang tidak membawa data mentah. Petunjuknya disertakan agar
            # pengguna tidak mencari-cari berkas yang tidak ada.
            arsip.writestr("sintaks/BACA_DULU.txt", _petunjuk_sintaks())
        for nomor, (judul, tabel) in enumerate(dokumen.tabel(), start=1):
            nama = "".join(
                c if c.isalnum() or c in "-_" else "_" for c in str(judul).lower()
            )[:40].strip("_")
            arsip.writestr(f"tabel/{nomor:02d}_{nama or 'tabel'}.csv", tabel.to_csv(index=False))
    return penampung.getvalue()


# --------------------------------------------------------------------------- #
# Daftar format
# --------------------------------------------------------------------------- #


@dataclass(frozen=True)
class Format:
    kode: str
    nama: str
    ekstensi: str
    mime: str
    keterangan: str


FORMAT: dict[str, Format] = {
    "docx": Format("docx", "Word", "docx",
                   "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                   "Dokumen siap disunting dan dicetak."),
    "pdf": Format("pdf", "PDF", "pdf", "application/pdf",
                  "Tata letak tetap, cocok untuk lampiran resmi."),
    "xlsx": Format("xlsx", "Excel", "xlsx",
                   "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                   "Tiap bagian menjadi lembar tersendiri; angka siap diolah ulang."),
    "pptx": Format("pptx", "PowerPoint", "pptx",
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                   "Slide siap dipresentasikan."),
    "html": Format("html", "HTML", "html", "text/html",
                   "Dibuka di peramban mana pun; dapat dicetak menjadi PDF."),
    "md": Format("md", "Markdown", "md", "text/markdown",
                 "Teks polos untuk disunting lebih lanjut."),
    "json": Format("json", "JSON", "json", "application/json",
                   "Data terstruktur untuk diolah sistem lain."),
    "py": Format("py", "Sintaks Python", "py", "text/x-python",
                 "Skrip yang menjalankan ulang analisis dengan pustaka yang sama."),
    "r": Format("r", "Sintaks R", "R", "text/plain",
                "Skrip R untuk memeriksa silang hasil di luar aplikasi."),
    "spss": Format("spss", "Sintaks SPSS", "sps", "text/plain",
                   "Perintah SPSS yang setara, untuk diperiksa ulang di sana."),
    "amos": Format("amos", "Spesifikasi AMOS", "txt", "text/plain",
                   "Daftar jalur dan langkah untuk dipindahkan ke AMOS."),
    "mplus": Format("mplus", "Input Mplus", "inp", "text/plain",
                    "Berkas .inp siap dijalankan di Mplus."),
    "zip": Format("zip", "Paket lengkap (ZIP)", "zip", "application/zip",
                  "Seluruh format, tabel dalam CSV, serta sintaks Python, R, SPSS, AMOS, dan Mplus."),
}

_PEMBUAT = {
    "docx": ke_docx,
    "pdf": ke_pdf,
    "xlsx": ke_xlsx,
    "pptx": ke_pptx,
    "html": ke_html,
    "md": ke_markdown,
    "json": ke_json,
    "py": ke_python,
    "r": ke_r,
    "spss": ke_spss,
    "amos": ke_amos,
    "mplus": ke_mplus,
    "zip": ke_zip,
}


def bangun(
    sumber, kode_format: str, pembaca: str = "eksekutif", lengkap: bool = False
) -> bytes:
    """Hasilkan berkas dalam format yang diminta.

    ``sumber`` boleh berupa ``Laporan`` naratif maupun ``Keranjang`` hasil yang
    dikumpulkan pengguna; keduanya diterjemahkan lebih dulu menjadi ``Dokumen``.
    """
    if kode_format not in _PEMBUAT:
        raise ValueError(
            f"Format '{kode_format}' tidak dikenal. Pilih dari: {', '.join(FORMAT)}."
        )
    if pembaca not in AUDIENCES:
        raise ValueError(f"Pembaca '{pembaca}' tidak dikenal.")
    return _PEMBUAT[kode_format](sumber, pembaca, lengkap)


def nama_berkas(
    sumber, kode_format: str, pembaca: str = "eksekutif", lengkap: bool = False
) -> str:
    if kode_format in {"py", "r", "spss", "amos", "mplus"}:
        ragam = "sintaks_analisis"
    elif isinstance(sumber, Keranjang):
        ragam = "laporan_hasil"
    else:
        ragam = (
            f"laporan_lengkap_{pembaca}" if lengkap else f"ringkasan_{pembaca}"
        )
    dasar = str(getattr(sumber, "dataset", "mvstatlab")).rsplit(".", 1)[0]
    bersih = "".join(c if c.isalnum() or c in "-_" else "_" for c in dasar)[:40].strip("_")
    return f"{bersih or 'lentera'}_{ragam}.{FORMAT[kode_format].ekstensi}"
